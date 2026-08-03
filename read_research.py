import sys
import os

try:
    from docx import Document
except ImportError:
    os.system('pip install python-docx')
    from docx import Document

try:
    from pptx import Presentation
except ImportError:
    os.system('pip install python-pptx')
    from pptx import Presentation

def read_docx(path):
    print(f"--- Reading {os.path.basename(path)} ---")
    doc = Document(path)
    text = []
    for p in doc.paragraphs:
        text.append(p.text)
    
    # Just print the first 5000 characters to avoid huge output, or all if short
    full = "\n".join(text)
    print(full[:10000])
    print("\n")

def read_pptx(path):
    print(f"--- Reading {os.path.basename(path)} ---")
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    full = "\n".join(text)
    print(full[:10000])
    print("\n")

read_docx(r"C:\Users\provu\Desktop\JRMSU LIBRARY LANDING PAGE\JRMSULibrary-web-PAGE\minutes for research 1.docx")
read_docx(r"C:\Users\provu\Desktop\JRMSU LIBRARY LANDING PAGE\JRMSULibrary-web-PAGE\JRMSU_FINAL_THESISHOPEFULLY.docx")
read_pptx(r"C:\Users\provu\Desktop\JRMSU LIBRARY LANDING PAGE\JRMSULibrary-web-PAGE\JRMSU_LMS_Proposal_FINAL_PLEASE.pptx")
